import os
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import filedialog, messagebox

def gerar_apresentacao(pasta, caminho_saida):
    imgs = sorted([f for f in os.listdir(pasta) if f.lower().endswith((".jpg", ".jpeg", ".png"))])
    if not imgs:
        messagebox.showerror("Erro", "Nenhuma imagem encontrada na pasta.")
        return

    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    blank_slide_layout = prs.slide_layouts[6]

    for img in imgs:
        slide = prs.slides.add_slide(blank_slide_layout)
        img_path = os.path.join(pasta, img)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    prs.save(caminho_saida)
    messagebox.showinfo("Sucesso", f"Apresentação salva com sucesso em:\n{caminho_saida}")

def escolher_pasta():
    pasta = filedialog.askdirectory(title="Selecione a pasta com as imagens")
    if not pasta:
        return

    # Pergunta onde salvar o arquivo
    caminho_saida = filedialog.asksaveasfilename(
        defaultextension=".pptx",
        filetypes=[("Apresentação PowerPoint", "*.pptx")],
        title="Salvar apresentação como...",
        initialfile="apresentacao_1920x1080.pptx"
    )
    if not caminho_saida:
        return

    gerar_apresentacao(pasta, caminho_saida)

# Interface
janela = tk.Tk()
janela.title("Gerador de Slides em PowerPoint")
janela.geometry("420x200")

label = tk.Label(janela, text="Clique no botão para selecionar a pasta com as imagens:")
label.pack(pady=20)

botao = tk.Button(janela, text="Selecionar Pasta e Gerar Apresentação", command=escolher_pasta)
botao.pack(pady=10)

janela.mainloop()
